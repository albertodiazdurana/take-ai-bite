#!/usr/bin/env python3
"""Convert non-markdown files to markdown for the Two-Pass Reading Protocol.

Supports: PDF, DOCX, HTML, PPTX
Output: Temporary markdown file in /tmp/

Usage:
    python scripts/convert_to_markdown.py <input_file>

The script prints the path to the converted markdown file on stdout.
Prefer system tools (pdftotext, pandoc) for speed; fall back to Python
libraries when system tools are unavailable.

Reference: DSM_0.2 Two-Pass Reading Strategy for Long Structured Files
"""

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


def convert_pdf(input_path: str, output_path: str) -> None:
    """Convert PDF to markdown. Tries pdftotext first, then pypdf."""
    if shutil.which("pdftotext"):
        subprocess.run(
            ["pdftotext", "-layout", input_path, output_path],
            check=True,
            capture_output=True,
        )
        return

    try:
        import pypdf
    except ImportError:
        raise RuntimeError(
            "No PDF converter available. Install pdftotext (poppler-utils) "
            "or pypdf (pip install pypdf)."
        )

    reader = pypdf.PdfReader(input_path)
    lines = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        lines.append(f"## Page {i}\n")
        lines.append(text.strip())
        lines.append("")
    Path(output_path).write_text("\n".join(lines), encoding="utf-8")


def convert_docx(input_path: str, output_path: str) -> None:
    """Convert DOCX to markdown. Tries pandoc first, then python-docx."""
    if shutil.which("pandoc"):
        subprocess.run(
            ["pandoc", "-f", "docx", "-t", "markdown", "-o", output_path, input_path],
            check=True,
            capture_output=True,
        )
        return

    try:
        import docx
    except ImportError:
        raise RuntimeError(
            "No DOCX converter available. Install pandoc or "
            "python-docx (pip install python-docx)."
        )

    doc = docx.Document(input_path)
    lines = []
    heading_map = {
        "Heading 1": "# ",
        "Heading 2": "## ",
        "Heading 3": "### ",
        "Heading 4": "#### ",
    }
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        prefix = heading_map.get(style_name, "")
        text = para.text.strip()
        if text:
            lines.append(f"{prefix}{text}")
            lines.append("")
    Path(output_path).write_text("\n".join(lines), encoding="utf-8")


def convert_html(input_path: str, output_path: str) -> None:
    """Convert HTML to markdown. Tries pandoc first, then markdownify."""
    if shutil.which("pandoc"):
        subprocess.run(
            ["pandoc", "-f", "html", "-t", "markdown", "-o", output_path, input_path],
            check=True,
            capture_output=True,
        )
        return

    try:
        import markdownify
    except ImportError:
        raise RuntimeError(
            "No HTML converter available. Install pandoc or "
            "markdownify (pip install markdownify)."
        )

    html_content = Path(input_path).read_text(encoding="utf-8")
    md_content = markdownify.markdownify(html_content, heading_style="ATX")
    Path(output_path).write_text(md_content, encoding="utf-8")


def convert_pptx(input_path: str, output_path: str) -> None:
    """Convert PPTX to markdown. Tries pandoc first, then python-pptx."""
    if shutil.which("pandoc"):
        subprocess.run(
            ["pandoc", "-f", "pptx", "-t", "markdown", "-o", output_path, input_path],
            check=True,
            capture_output=True,
        )
        return

    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError(
            "No PPTX converter available. Install pandoc or "
            "python-pptx (pip install python-pptx)."
        )

    prs = Presentation(input_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append("")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
                lines.append("")
    Path(output_path).write_text("\n".join(lines), encoding="utf-8")


CONVERTERS = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".html": convert_html,
    ".htm": convert_html,
    ".pptx": convert_pptx,
}


def main() -> None:
    if len(sys.argv) != 2:
        print(f"Usage: {sys.argv[0]} <input_file>", file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]
    if not os.path.isfile(input_path):
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    ext = Path(input_path).suffix.lower()
    if ext not in CONVERTERS:
        supported = ", ".join(sorted(CONVERTERS.keys()))
        print(
            f"Error: unsupported format '{ext}'. Supported: {supported}",
            file=sys.stderr,
        )
        sys.exit(1)

    stem = Path(input_path).stem
    fd, output_path = tempfile.mkstemp(prefix=f"{stem}_", suffix=".md", dir="/tmp")
    os.close(fd)

    try:
        CONVERTERS[ext](input_path, output_path)
    except Exception as e:
        os.unlink(output_path)
        print(f"Error: conversion failed: {e}", file=sys.stderr)
        sys.exit(1)

    print(output_path)


if __name__ == "__main__":
    main()
